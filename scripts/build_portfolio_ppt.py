from pathlib import Path
import shutil
import argparse

import fitz
from pptx import Presentation
from pptx.util import Inches


ROOT = Path(__file__).resolve().parents[1]
RENDER_DIR = ROOT / ".ppt-render"


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("source", type=Path)
    parser.add_argument("output", type=Path)
    args = parser.parse_args()
    source = args.source.resolve()
    output = args.output.resolve()

    if RENDER_DIR.exists():
        shutil.rmtree(RENDER_DIR)
    RENDER_DIR.mkdir()

    document = fitz.open(source)
    first = document[0].rect

    presentation = Presentation()
    presentation.slide_width = Inches(7.5)
    presentation.slide_height = int(presentation.slide_width * first.height / first.width)
    blank = presentation.slide_layouts[6]

    for index, page in enumerate(document):
        image_path = RENDER_DIR / f"page-{index + 1:02d}.png"
        pixmap = page.get_pixmap(matrix=fitz.Matrix(2.4, 2.4), alpha=False)
        pixmap.save(image_path)

        slide = presentation.slides.add_slide(blank)
        slide.shapes.add_picture(
            str(image_path), 0, 0,
            width=presentation.slide_width,
            height=presentation.slide_height,
        )

        # Preserve clickable links from the source PDF with transparent overlays.
        scale_x = presentation.slide_width / page.rect.width
        scale_y = presentation.slide_height / page.rect.height
        for link in page.get_links():
            uri = link.get("uri")
            rect = link.get("from")
            if not uri or not rect:
                continue
            overlay = slide.shapes.add_shape(
                1,
                int(rect.x0 * scale_x),
                int(rect.y0 * scale_y),
                max(1, int(rect.width * scale_x)),
                max(1, int(rect.height * scale_y)),
            )
            overlay.fill.background()
            overlay.line.fill.background()
            overlay.click_action.hyperlink.address = uri

    presentation.core_properties.title = "김도혁 포트폴리오 2026"
    presentation.core_properties.subject = "Web Frontend Developer Portfolio"
    presentation.core_properties.author = "김도혁"
    presentation.core_properties.comments = "기존 포트폴리오 PDF를 고해상도 PPT로 변환"
    output.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(output)
    shutil.rmtree(RENDER_DIR)
    print(f"created={output}")
    print(f"slides={len(presentation.slides)}")


if __name__ == "__main__":
    main()
